"""PPTX ビルダー（python-pptx でスライドを組み立て）"""
import json
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from src.config import (
    SLIDE_WIDTH_INCHES,
    SLIDE_HEIGHT_INCHES,
    DEFAULT_FONT,
)

# テキスト色のマッピング
TEXT_COLORS = {
    "dark": RGBColor(255, 255, 255),    # 白文字
    "light": RGBColor(33, 33, 33),       # 暗い文字
    "colorful": RGBColor(255, 255, 255), # 白文字
}

# オーバーレイ色のマッピング
OVERLAY_COLORS = {
    "dark": (RGBColor(0, 0, 0), "40000"),       # 黒 40%
    "light": (RGBColor(255, 255, 255), "50000"), # 白 50%
    "colorful": (RGBColor(0, 0, 0), "30000"),    # 黒 30%
}


def build_pptx(
    plan: dict,
    image_paths: list[Path | None],
    output_dir: Path,
    taste: dict,
) -> Path:
    """
    スライド構成 JSON + 背景画像からPPTXを組み立てる。

    Returns:
        生成された PPTX のファイルパス
    """
    output_dir = Path(output_dir)
    session_id = f"session_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    session_dir = output_dir / session_id
    session_dir.mkdir(parents=True, exist_ok=True)

    # 構成 JSON を保存（リライト用）
    plan_path = session_dir / "plan.json"
    plan_path.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    color_tone = taste.get("color_tone", "dark")
    text_color = TEXT_COLORS.get(color_tone, TEXT_COLORS["dark"])

    for i, slide_data in enumerate(plan["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト

        # 背景画像を配置
        img_path = image_paths[i] if i < len(image_paths) else None
        if img_path and Path(img_path).exists():
            _add_background_image(slide, prs, img_path)
            _add_overlay(slide, prs, color_tone)

        # テキスト配置
        slide_type = slide_data.get("slide_type", "content")
        _layout_text(slide, prs, slide_data, slide_type, text_color)

    pptx_path = session_dir / "presentation.pptx"
    prs.save(str(pptx_path))

    # 画像をセッションディレクトリにコピー（リライト用）
    images_dest = session_dir / "images"
    images_dest.mkdir(exist_ok=True)
    for img_path in image_paths:
        if img_path and Path(img_path).exists():
            import shutil
            shutil.copy2(img_path, images_dest / Path(img_path).name)

    return pptx_path


def _add_background_image(slide, prs, image_path: Path) -> None:
    """スライド全体に背景画像を配置し、最背面に移動"""
    pic = slide.shapes.add_picture(
        str(image_path),
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    # 最背面に移動（XML操作）
    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)


def _add_overlay(slide, prs, color_tone: str) -> None:
    """半透明オーバーレイを追加（テキスト可読性向上）"""
    overlay_color, alpha_val = OVERLAY_COLORS.get(
        color_tone, OVERLAY_COLORS["dark"]
    )

    overlay = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = overlay_color

    # 透明度の設定（XML操作）
    solid_fill = overlay.fill._fill
    srgb_clr = solid_fill.find(qn("a:srgbClr"))
    if srgb_clr is not None:
        alpha = srgb_clr.makeelement(qn("a:alpha"), {"val": alpha_val})
        srgb_clr.append(alpha)

    # 枠線を消す
    overlay.line.fill.background()


def _layout_text(slide, prs, slide_data: dict, slide_type: str, text_color: RGBColor) -> None:
    """スライドタイプに応じたテキスト配置"""
    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    bullets = slide_data.get("bullet_points", [])
    notes = slide_data.get("speaker_notes", "")

    if slide_type == "title":
        _layout_title_slide(slide, prs, title, subtitle, text_color)
    elif slide_type == "section":
        _layout_section_slide(slide, prs, title, text_color)
    elif slide_type == "end":
        _layout_end_slide(slide, prs, title, text_color)
    else:
        _layout_content_slide(slide, prs, title, bullets, text_color)

    # スピーカーノートを追加
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _layout_title_slide(slide, prs, title: str, subtitle: str, color: RGBColor) -> None:
    """表紙スライド"""
    # タイトル
    _add_textbox(
        slide,
        text=title,
        left=Inches(1.5),
        top=Inches(2.5),
        width=Inches(10),
        height=Inches(1.5),
        font_size=Pt(44),
        font_color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    # サブタイトル
    if subtitle:
        _add_textbox(
            slide,
            text=subtitle,
            left=Inches(2),
            top=Inches(4.2),
            width=Inches(9),
            height=Inches(1),
            font_size=Pt(24),
            font_color=color,
            alignment=PP_ALIGN.CENTER,
        )


def _layout_section_slide(slide, prs, title: str, color: RGBColor) -> None:
    """セクション区切りスライド"""
    _add_textbox(
        slide,
        text=title,
        left=Inches(1.5),
        top=Inches(3),
        width=Inches(10),
        height=Inches(1.5),
        font_size=Pt(40),
        font_color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def _layout_content_slide(slide, prs, title: str, bullets: list[str], color: RGBColor) -> None:
    """コンテンツスライド"""
    # タイトル
    _add_textbox(
        slide,
        text=title,
        left=Inches(0.8),
        top=Inches(0.4),
        width=Inches(11),
        height=Inches(0.8),
        font_size=Pt(32),
        font_color=color,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    # 箇条書き
    if bullets:
        txbox = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(1.5),
            width=Inches(10),
            height=Inches(5),
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        for idx, bullet in enumerate(bullets):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.name = DEFAULT_FONT
            p.font.color.rgb = color
            p.space_after = Pt(8)
            p.alignment = PP_ALIGN.LEFT


def _layout_end_slide(slide, prs, title: str, color: RGBColor) -> None:
    """最終スライド"""
    _add_textbox(
        slide,
        text=title,
        left=Inches(1.5),
        top=Inches(3),
        width=Inches(10),
        height=Inches(1.5),
        font_size=Pt(36),
        font_color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def _add_textbox(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    font_size=Pt(20),
    font_color=RGBColor(255, 255, 255),
    bold=False,
    alignment=PP_ALIGN.LEFT,
) -> None:
    """テキストボックスを追加するヘルパー"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = DEFAULT_FONT
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
